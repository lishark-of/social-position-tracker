from __future__ import annotations

import io
import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt


ROOT = Path("/Users/shark-li/Documents/New project")
SRC = ROOT / "鞣花酸-断奶仔猪文献汇报_逐图拆解_博士汇报优化版_含讲稿备注.pptx"
OUT = ROOT / "鞣花酸-断奶仔猪文献汇报_图像精修版_博士汇报优化版_含讲稿备注.pptx"
WORK = ROOT / "outputs/manual-ellagic/presentations/ellagic-piglet"
EXTRACTED = WORK / "assets/extracted"
ASSET_OUT = WORK / "assets/refined_panels"

EMU_PER_IN = 914400
DARK = RGBColor(36, 43, 48)
RED = RGBColor(168, 48, 42)
MUTED = RGBColor(95, 105, 112)


def in_(x: float) -> int:
    return int(x * EMU_PER_IN)


def crop(name: str, box: tuple[int, int, int, int], out_name: str, border: int = 0) -> Path:
    im = Image.open(EXTRACTED / name).convert("RGB")
    c = im.crop(box)
    if border:
        c = ImageOps.expand(c, border=border, fill="white")
    path = ASSET_OUT / out_name
    c.save(path, quality=95)
    return path


def crop_from_slide(prs: Presentation, slide_no: int, out_name: str, box: tuple[int, int, int, int], pic_index: int = 0) -> Path:
    slide = prs.slides[slide_no - 1]
    pics = [shp for shp in slide.shapes if shp.shape_type == 13]
    pic = pics[pic_index]
    im = Image.open(io.BytesIO(pic.image.blob)).convert("RGB")
    path = ASSET_OUT / out_name
    im.crop(box).save(path, quality=95)
    return path


def crop_fig2_bd_montage() -> Path:
    im = Image.open(EXTRACTED / "animal_p05_img2_1718x980.png").convert("RGB")
    b = im.crop((1005, 0, 1368, 492))
    c = im.crop((1370, 0, 1718, 492))
    d = im.crop((1005, 492, 1368, 980))
    canvas = Image.new("RGB", (735, 985), "white")
    canvas.paste(b, (0, 0))
    canvas.paste(c, (375, 0))
    canvas.paste(d, (185, 500))
    path = ASSET_OUT / "animal_fig2BD_clean.png"
    ImageOps.expand(canvas, border=8, fill="white").save(path, quality=95)
    return path


def remove_pictures(slide):
    for shp in list(slide.shapes):
        if shp.shape_type == 13:
            el = shp._element
            el.getparent().remove(el)


def place(slide, image_path: Path, x: float, y: float, w: float, h: float):
    im = Image.open(image_path)
    ratio = im.width / im.height
    box_ratio = w / h
    if ratio >= box_ratio:
        width = w
        height = w / ratio
    else:
        height = h
        width = h * ratio
    left = x + (w - width) / 2
    top = y + (h - height) / 2
    return slide.shapes.add_picture(str(image_path), in_(left), in_(top), width=in_(width), height=in_(height))


def set_text(shape, text: str, size=16, bold=False, color=DARK):
    shape.text = text
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def text_shapes(slide):
    return [s for s in slide.shapes if getattr(s, "has_text_frame", False) and s.text.strip()]


def update_body(slide, heading: str, lines: list[str]):
    shapes = text_shapes(slide)
    # Use right-side heading/body boxes in the existing template.
    if len(shapes) >= 2:
        candidates = [s for s in shapes if s.left > in_(8.0) and s.top > in_(0.8) and not s.text.strip().isdigit()]
        if len(candidates) >= 2:
            set_text(candidates[0], heading, size=16, bold=True, color=RED)
            set_text(candidates[1], "\n".join(lines), size=16, color=DARK)


def add_source(slide, text: str):
    # Remove previous source tags from earlier runs.
    for shp in list(slide.shapes):
        if getattr(shp, "has_text_frame", False) and shp.text.strip().startswith("Source:"):
            el = shp._element
            el.getparent().remove(el)
    tb = slide.shapes.add_textbox(in_(0.72), in_(6.88), in_(5.9), in_(0.18))
    set_text(tb, text, size=7, color=MUTED)


def replace_slide(slide, images: list[Path], layout: str = "left", source: str = "Source: Animal Nutrition, 2022"):
    remove_pictures(slide)
    if layout == "left":
        place(slide, images[0], 0.72, 0.96, 7.55, 5.78)
    elif layout == "left2":
        place(slide, images[0], 0.70, 0.98, 3.75, 5.62)
        place(slide, images[1], 4.55, 0.98, 3.75, 5.62)
    elif layout == "wide":
        place(slide, images[0], 0.72, 1.02, 7.75, 5.35)
    elif layout == "top2":
        place(slide, images[0], 0.72, 0.96, 3.75, 5.72)
        place(slide, images[1], 4.50, 0.96, 3.85, 5.72)
    elif layout == "summary2":
        place(slide, images[0], 0.72, 1.0, 3.75, 5.55)
        place(slide, images[1], 4.70, 1.0, 3.65, 5.55)
    add_source(slide, source)


def first_title(slide):
    candidates = []
    for shp in text_shapes(slide):
        txt = shp.text.strip()
        if txt.startswith("主讲：") or txt.startswith("Source:") or txt.isdigit():
            continue
        candidates.append(shp)
    if not candidates:
        return ""
    shp = sorted(candidates, key=lambda s: (s.top, s.left))[0]
    return shp.text.strip().replace("\n", " | ")


def find_slide(prs: Presentation, prefix: str):
    for i, slide in enumerate(prs.slides, 1):
        if first_title(slide).startswith(prefix):
            return i, slide
    raise KeyError(prefix)


def note_check(pptx_path: Path) -> tuple[int, list[int], list[int]]:
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main", "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    with zipfile.ZipFile(pptx_path) as z:
        pres = ET.fromstring(z.read("ppt/presentation.xml"))
        rels = ET.fromstring(z.read("ppt/_rels/presentation.xml.rels"))
        rid_to_target = {el.attrib["Id"]: el.attrib["Target"] for el in rels}
        missing, bad = [], []
        count = 0
        for idx, sld in enumerate(pres.find("p:sldIdLst", ns), 1):
            rid = sld.attrib["{" + ns["r"] + "}id"]
            slide = "ppt/" + rid_to_target[rid]
            rel_path = "ppt/slides/_rels/" + slide.split("/")[-1] + ".rels"
            root = ET.fromstring(z.read(rel_path))
            note = None
            for el in root.findall("{%s}Relationship" % rel_ns):
                if el.attrib.get("Type", "").endswith("/notesSlide"):
                    note = "ppt/notesSlides/" + el.attrib["Target"].split("/")[-1]
            if not note or note not in z.namelist():
                missing.append(idx)
                continue
            count += 1
            text = re.sub("<[^>]+>", " ", z.read(note).decode("utf-8", "ignore"))
            if not all(tok in text for tok in ["讲解逻辑：", "本页定位", "图怎么看", "结果判断", "机制解释", "过渡句"]):
                bad.append(idx)
        return count, missing, bad


def main():
    ASSET_OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation(SRC)

    # Animal Nutrition: direct embedded figure crops.
    panels = {
        "table3_clean": crop_from_slide(prs, 7, "table3_clean.png", (0, 155, 795, 510)),
        "fig1A": crop("animal_p05_img1_1600x1209.png", (0, 0, 790, 715), "animal_fig1A_clean.png", 8),
        "fig1BCD": crop("animal_p05_img1_1600x1209.png", (850, 0, 1600, 715), "animal_fig1BCD_clean.png", 8),
        "fig1EFG": crop("animal_p05_img1_1600x1209.png", (0, 735, 900, 1209), "animal_fig1EFG_clean.png", 8),
        "fig1HI": crop("animal_p05_img1_1600x1209.png", (900, 735, 1600, 1209), "animal_fig1HI_clean.png", 8),
        "fig2A": crop("animal_p05_img2_1718x980.png", (0, 0, 1005, 980), "animal_fig2A_clean.png", 8),
        "fig2BD": crop_fig2_bd_montage(),
        "fig2E": crop("animal_p05_img2_1718x980.png", (1370, 485, 1718, 980), "animal_fig2E_clean.png", 8),
        "fig3AB": crop("animal_p06_img1_1836x1097.png", (0, 0, 920, 575), "animal_fig3AB_clean.png", 8),
        "fig3CD": crop("animal_p06_img1_1836x1097.png", (940, 0, 1836, 575), "animal_fig3CD_clean.png", 8),
        "fig3EG": crop("animal_p06_img1_1836x1097.png", (0, 590, 1115, 1097), "animal_fig3EG_clean.png", 8),
        "fig3HI": crop("animal_p06_img1_1836x1097.png", (1115, 590, 1836, 1097), "animal_fig3HI_clean.png", 8),
        "fig4AB": crop("animal_p07_img1_1588x1605.png", (0, 0, 1588, 710), "animal_fig4AB_clean.png", 8),
        "fig4C": crop("animal_p07_img1_1588x1605.png", (0, 760, 1588, 1605), "animal_fig4C_clean.png", 8),
        "table4_clean": crop_from_slide(prs, 26, "table4_clean.png", (0, 0, 795, 238), pic_index=1),
        "fig5AB": crop("animal_p08_img1_1836x2096.png", (0, 0, 1836, 980), "animal_fig5AB_clean.png", 8),
        "fig5CD": crop("animal_p08_img1_1836x2096.png", (0, 1000, 870, 1600), "animal_fig5CD_clean.png", 8),
        "fig5EF": crop("animal_p08_img1_1836x2096.png", (870, 1000, 1836, 1600), "animal_fig5EF_clean.png", 8),
        "fig5GH": crop("animal_p08_img1_1836x2096.png", (0, 1600, 900, 2096), "animal_fig5GH_clean.png", 8),
        "fig5IJ": crop("animal_p08_img1_1836x2096.png", (880, 1600, 1836, 2096), "animal_fig5IJ_clean.png", 8),
        "fig6A": crop("animal_p09_img1_1718x1473.png", (0, 0, 1718, 780), "animal_fig6A_clean.png", 8),
        "fig6B": crop("animal_p09_img1_1718x1473.png", (0, 790, 1718, 1473), "animal_fig6B_clean.png", 8),
        "table5_clean": crop_from_slide(prs, 38, "table5_clean.png", (0, 82, 810, 548)),
        "fig7": crop("animal_p10_img1_1541x725.png", (0, 0, 1541, 725), "animal_fig7_clean.png", 8),
    }

    # Microbiome: direct embedded figure crops.
    panels.update(
        {
            "micro2abc": crop("micro_p09_img1_1709x2481.png", (0, 0, 1709, 480), "micro_fig2abc_clean.png", 8),
            "micro2d": crop("micro_p09_img1_1709x2481.png", (0, 485, 1709, 690), "micro_fig2d_clean.png", 8),
            "micro2eg": crop("micro_p09_img1_1709x2481.png", (0, 700, 1709, 1148), "micro_fig2eg_clean.png", 8),
            "micro2h": crop("micro_p09_img1_1709x2481.png", (0, 1160, 610, 1555), "micro_fig2h_clean.png", 8),
            "micro2ik": crop("micro_p09_img1_1709x2481.png", (610, 1160, 1709, 1555), "micro_fig2ik_clean.png", 8),
            "micro2ln": crop("micro_p09_img1_1709x2481.png", (0, 1560, 1709, 2015), "micro_fig2ln_clean.png", 8),
            "micro2or": crop("micro_p09_img1_1709x2481.png", (0, 2015, 1709, 2481), "micro_fig2or_clean.png", 8),
            "micro7ac": crop("micro_p18_img1_1718x2481.png", (0, 0, 1718, 430), "micro_fig7ac_clean.png", 8),
            "micro7dg": crop("micro_p18_img1_1718x2481.png", (0, 430, 1718, 925), "micro_fig7dg_clean.png", 8),
            "micro7h": crop("micro_p18_img1_1718x2481.png", (0, 950, 610, 1370), "micro_fig7h_clean.png", 8),
            "micro7ik": crop("micro_p18_img1_1718x2481.png", (610, 950, 1718, 1370), "micro_fig7ik_clean.png", 8),
            "micro7lm": crop("micro_p18_img1_1718x2481.png", (0, 1380, 1718, 1945), "micro_fig7lm_clean.png", 8),
            "micro7no": crop("micro_p18_img1_1718x2481.png", (0, 1950, 1718, 2481), "micro_fig7no_clean.png", 8),
            "micro8ad": crop("micro_p20_img1_1755x2481.png", (0, 0, 1755, 555), "micro_fig8ad_clean.png", 8),
            "micro8eh": crop("micro_p20_img1_1755x2481.png", (0, 565, 1755, 1070), "micro_fig8eh_clean.png", 8),
            "micro8io": crop("micro_p20_img1_1755x2481.png", (0, 1085, 1755, 2010), "micro_fig8io_clean.png", 8),
            "micro8ps": crop("micro_p20_img1_1755x2481.png", (0, 2030, 1755, 2481), "micro_fig8ps_clean.png", 8),
        }
    )

    replacements = {
        "Table 3": ([panels["table3_clean"]], "left"),
        "Figure 1A": ([panels["fig1A"]], "left"),
        "Figure 1B-D": ([panels["fig1BCD"]], "left"),
        "Figure 1E-G": ([panels["fig1EFG"]], "left"),
        "Figure 1H-I": ([panels["fig1HI"]], "left"),
        "Figure 1 小结": ([panels["fig1A"], panels["fig1HI"]], "left2"),
        "Figure 2A": ([panels["fig2A"]], "left"),
        "Figure 2B-D": ([panels["fig2BD"]], "left"),
        "Figure 2E": ([panels["fig2E"]], "left"),
        "Figure 2 小结": ([panels["fig2A"], panels["fig2BD"]], "left2"),
        "Figure 3A-B": ([panels["fig3AB"]], "left"),
        "Figure 3C-D": ([panels["fig3CD"]], "left"),
        "Figure 3E-G": ([panels["fig3EG"]], "left"),
        "Figure 3H-I": ([panels["fig3HI"]], "left"),
        "Figure 3 小结": ([panels["fig3AB"], panels["fig3CD"]], "left2"),
        "Figure 4A-B": ([panels["fig4AB"]], "left"),
        "Figure 4C/Table 4": ([panels["fig4C"], panels["table4_clean"]], "left2"),
        "Figure 4 小结": ([panels["fig4AB"], panels["fig4C"]], "left2"),
        "Figure 5A-B": ([panels["fig5AB"]], "left"),
        "Figure 5C-D": ([panels["fig5CD"]], "left"),
        "Figure 5E-F": ([panels["fig5EF"]], "left"),
        "Figure 5G-H": ([panels["fig5GH"]], "left"),
        "Figure 5I-J": ([panels["fig5IJ"]], "left"),
        "Figure 5 小结": ([panels["fig5EF"], panels["fig5GH"]], "left2"),
        "Figure 6A": ([panels["fig6A"]], "left"),
        "Figure 6B": ([panels["fig6B"]], "left"),
        "Table 5": ([panels["table5_clean"]], "left"),
        "Figure 6 小结": ([panels["fig6A"], panels["table5_clean"]], "left2"),
        "Figure 7": ([panels["fig7"]], "left"),
        "Microbiome Fig.2a-d": ([panels["micro2abc"], panels["micro2d"]], "left2"),
        "Microbiome Fig.2e-g": ([panels["micro2eg"]], "left"),
        "Microbiome Fig.2h-k": ([panels["micro2h"], panels["micro2ik"]], "left2"),
        "Microbiome Fig.2l-r": ([panels["micro2ln"], panels["micro2or"]], "left2"),
        "Microbiome Fig.7a-g": ([panels["micro7ac"], panels["micro7dg"]], "left2"),
        "Microbiome Fig.7h-k": ([panels["micro7h"], panels["micro7ik"]], "left2"),
        "Microbiome Fig.7l-o": ([panels["micro7lm"], panels["micro7no"]], "left2"),
        "Microbiome Fig.8a-h": ([panels["micro8ad"], panels["micro8eh"]], "left2"),
        "Microbiome Fig.8i-o": ([panels["micro8io"]], "left"),
        "Microbiome Fig.8p-s": ([panels["micro8ps"]], "left"),
    }

    animal_source = "Source: Animal Nutrition, 2022"
    micro_source = "Source: Microbiome, 2023"
    touched = []
    for prefix, (imgs, layout) in replacements.items():
        try:
            slide_no, slide = find_slide(prs, prefix)
        except KeyError:
            continue
        source = micro_source if prefix.startswith("Microbiome") else animal_source
        replace_slide(slide, imgs, layout, source)
        touched.append(slide_no)

    # Tighten result boxes for the worst affected pages.
    body_updates = {
        "Figure 3E-G": (
            "结果要点",
            [
                "1. 图怎么看：检测 TNF-α、IL-6、IL-1α mRNA。",
                "2. 主要结果：EA/FEA 对多数促炎因子影响有限。",
                "3. 机制含义：炎症不是本文最稳定主线。",
                "4. 本页结论：后续重点转向菌群和代谢物。",
            ],
        ),
        "Figure 3H-I": (
            "结果要点",
            [
                "1. 图怎么看：继续观察 IL-1β 和 IL-10。",
                "2. 主要结果：FEA 对 IL-10 有上升趋势。",
                "3. 机制含义：提示温和调节，不是核心效应。",
                "4. 本页结论：菌群结构结果更关键。",
            ],
        ),
        "整合机制假说": (
            "汇报时必须区分",
            [
                "主文献已验证：EA 改善腹泻、氧化应激、屏障、菌群和 SCFAs。",
                "机制补充：3-PPA 可通过 AhR 增强猪肠屏障。",
                "本人课题待验证：荔枝酚类是否通过芳香族代谢物/SCFAs—AhR/Nrf2 修复屏障。",
            ],
        ),
    }
    for prefix, (head, lines) in body_updates.items():
        try:
            _, slide = find_slide(prs, prefix)
            update_body(slide, head, lines)
        except KeyError:
            pass

    prs.save(OUT)
    notes_count, missing_notes, bad_notes = note_check(OUT)

    # Simple post-save picture density and source-asset check.
    check = Presentation(OUT)
    over = []
    for idx, slide in enumerate(check.slides, 1):
        pics = sum(1 for shp in slide.shapes if shp.shape_type == 13)
        if pics > 3:
            over.append((idx, pics))
    print(OUT)
    print("slides", len(check.slides))
    print("touched", sorted(touched))
    print("notes", notes_count, "missing", missing_notes, "bad", bad_notes)
    print("over3", over)


if __name__ == "__main__":
    main()
