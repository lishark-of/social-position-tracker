from __future__ import annotations

import io
import re
import shutil
import zipfile
from html import escape
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path("/Users/shark-li/Documents/New project")
SRC = ROOT / "鞣花酸-断奶仔猪文献汇报_逐图拆解_优化排版_含讲稿备注.pptx"
WORK = ROOT / "outputs/manual-ellagic/presentations/ellagic-piglet"
ASSETS = WORK / "assets/optimized"
OUT = ROOT / "鞣花酸-断奶仔猪文献汇报_逐图拆解_博士汇报优化版_含讲稿备注.pptx"
INTERMEDIATE = WORK / "output/ellagic_doctor_optimized_intermediate.pptx"

EMU_PER_IN = 914400
SLIDE_W = 13.333
SLIDE_H = 7.5

RED = RGBColor(168, 48, 42)
DARK = RGBColor(36, 43, 48)
MUTED = RGBColor(95, 105, 112)
PALE = RGBColor(248, 243, 239)
LINE = RGBColor(222, 220, 216)


def in_(x: float) -> int:
    return int(x * EMU_PER_IN)


def text_shapes(slide):
    return [s for s in slide.shapes if getattr(s, "has_text_frame", False) and s.text.strip()]


def title_shape(slide):
    shapes = text_shapes(slide)
    return shapes[0] if shapes else None


def set_text(shape, text: str, size=18, bold=False, color=DARK, font="Microsoft YaHei"):
    shape.text = text
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def set_title(slide, text: str):
    shp = title_shape(slide)
    if shp is not None:
        set_text(shp, text, size=24, bold=True, color=DARK)


def set_tag(slide, text: str):
    shapes = text_shapes(slide)
    if len(shapes) > 1:
        set_text(shapes[1], text, size=14, bold=True, color=RED)


def update_body(slide, heading: str, body: str):
    shapes = text_shapes(slide)
    if len(shapes) >= 6:
        set_text(shapes[-2], heading, size=16, bold=True, color=RED)
        set_text(shapes[-1], body, size=17, color=DARK)


def remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def remove_pictures(slide):
    for shp in list(slide.shapes):
        if shp.shape_type == 13:
            remove_shape(shp)


def remove_content_except_footer(slide):
    for shp in list(slide.shapes):
        keep = False
        if getattr(shp, "has_text_frame", False):
            txt = shp.text.strip()
            if txt.startswith("主讲：") or txt.isdigit():
                keep = True
        if not keep:
            remove_shape(shp)


def add_footer(slide, number: int | None = None):
    footer = slide.shapes.add_textbox(in_(0.45), in_(7.17), in_(9.6), in_(0.18))
    set_text(footer, "主讲：Animal Nutrition（鞣花酸-断奶仔猪）｜机制补充：Microbiome（3-PPA/AhR）", size=8, color=MUTED)
    num = slide.shapes.add_textbox(in_(12.36), in_(7.12), in_(0.45), in_(0.22))
    set_text(num, str(number or ""), size=10, color=MUTED)


def add_title_bar(slide, title: str, tag: str):
    title_box = slide.shapes.add_textbox(in_(0.45), in_(0.18), in_(10.4), in_(0.48))
    set_text(title_box, title, size=24, bold=True, color=DARK)
    tag_box = slide.shapes.add_textbox(in_(11.0), in_(0.22), in_(1.85), in_(0.35))
    set_text(tag_box, tag, size=14, bold=True, color=RED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, in_(0.45), in_(0.76), in_(12.4), in_(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.color.rgb = LINE


def add_text_block(slide, x, y, w, h, title: str, lines: list[str], font_size=16):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, in_(x), in_(y), in_(w), in_(h))
    box.fill.solid()
    box.fill.fore_color.rgb = PALE
    box.line.color.rgb = LINE
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.runs[0].font.name = "Microsoft YaHei"
    p.runs[0].font.size = Pt(font_size + 1)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RED
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        for r in p.runs:
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(font_size)
            r.font.color.rgb = DARK
    return box


def add_content_slide(prs, title: str, tag: str, blocks: list[tuple[float, float, float, float, str, list[str], int]]):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_title_bar(slide, title, tag)
    for block in blocks:
        add_text_block(slide, *block)
    add_footer(slide)
    return slide


def add_figure_slide(prs, title: str, tag: str, img_paths: list[Path], body_heading: str, body: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_title_bar(slide, title, tag)
    if len(img_paths) == 1:
        slide.shapes.add_picture(str(img_paths[0]), in_(1.1), in_(1.05), width=in_(6.9), height=in_(5.2))
    else:
        slide.shapes.add_picture(str(img_paths[0]), in_(0.75), in_(1.25), width=in_(3.7))
        slide.shapes.add_picture(str(img_paths[1]), in_(4.65), in_(1.25), width=in_(3.75))
    head = slide.shapes.add_textbox(in_(8.95), in_(1.12), in_(3.5), in_(0.35))
    set_text(head, body_heading, size=16, bold=True, color=RED)
    txt = slide.shapes.add_textbox(in_(8.95), in_(1.55), in_(3.55), in_(4.9))
    set_text(txt, body, size=17, color=DARK)
    add_footer(slide)
    return slide


def move_last_slide_to(prs, index_zero_based: int):
    sld_id_lst = prs.slides._sldIdLst
    last = sld_id_lst[-1]
    sld_id_lst.remove(last)
    sld_id_lst.insert(index_zero_based, last)


def renumber(slide_collection):
    for idx, slide in enumerate(slide_collection, 1):
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False) and shp.text.strip().isdigit():
                if shp.left > in_(11.5) and shp.top > in_(6.7):
                    set_text(shp, str(idx), size=10, color=MUTED)


def first_title(slide):
    candidates = []
    for shp in text_shapes(slide):
        txt = shp.text.strip()
        if txt.startswith("主讲：") or txt.isdigit():
            continue
        candidates.append(shp)
    if not candidates:
        return ""
    shp = sorted(candidates, key=lambda s: (s.top, s.left))[0]
    return shp.text.strip().replace("\n", " | ")


def crop_fig2_panels():
    src = WORK / "assets/extracted/animal_p05_img2_1718x980.png"
    im = Image.open(src).convert("RGB")
    b = im.crop((1005, 0, 1368, 492))
    c = im.crop((1370, 0, 1718, 492))
    d = im.crop((1005, 492, 1368, 980))
    e = im.crop((1370, 492, 1718, 980))
    canvas = Image.new("RGB", (735, 985), "white")
    canvas.paste(b, (0, 0))
    canvas.paste(c, (375, 0))
    canvas.paste(d, (185, 500))
    bd = ASSETS / "fig2_b_d_morphometry.png"
    ep = ASSETS / "fig2e_tunel_aod.png"
    canvas.save(bd, quality=95)
    e.save(ep, quality=95)
    return bd, ep


def patch_note_text_xml(xml: str, note_lines: list[str], slide_num: int) -> str:
    body = ["讲解逻辑："] + [f"{i}. {line}" for i, line in enumerate(note_lines, 1)]
    tx = ['<p:txBody><a:bodyPr/><a:lstStyle/>']
    for para in body:
        tx.append(
            '<a:p><a:r><a:rPr lang="zh-CN" dirty="0"/>'
            f"<a:t>{escape(para)}</a:t></a:r><a:endParaRPr lang=\"zh-CN\" dirty=\"0\"/></a:p>"
        )
    tx.append("</p:txBody>")
    marker = "Notes Placeholder 2"
    if marker in xml:
        pos = xml.index(marker)
        start = xml.index("<p:txBody>", pos)
        end = xml.index("</p:txBody>", start) + len("</p:txBody>")
        xml = xml[:start] + "".join(tx) + xml[end:]
    else:
        xml = re.sub(r"<p:txBody>.*?</p:txBody>", "".join(tx), xml, count=1, flags=re.S)
    xml = re.sub(r"<a:t>\d+</a:t>", f"<a:t>{slide_num}</a:t>", xml, count=1)
    return xml


def make_note_xml(note_lines: list[str], slide_num: int) -> str:
    base = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr><p:sp><p:nvSpPr><p:cNvPr id="2" name="Slide Image Placeholder 1"/><p:cNvSpPr><a:spLocks noGrp="1" noRot="1" noChangeAspect="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldImg"/></p:nvPr></p:nvSpPr><p:spPr/></p:sp><p:sp><p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder 2"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="zh-CN" dirty="0"/><a:t></a:t></a:r><a:endParaRPr lang="zh-CN" dirty="0"/></a:p></p:txBody></p:sp><p:sp><p:nvSpPr><p:cNvPr id="4" name="Slide Number Placeholder 3"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldNum" sz="quarter" idx="10"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:fld id="{F7021451-1387-4CA6-816F-3879F97B5CBC}" type="slidenum"><a:rPr lang="en-US"/><a:t>1</a:t></a:fld><a:endParaRPr lang="en-US"/></a:p></p:txBody></p:sp></p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:notes>"""
    return patch_note_text_xml(base, note_lines, slide_num)


def note_lines(title: str) -> list[str]:
    t = title
    if "封面" in t or t.startswith("鞣花酸"):
        return [
            "本页定位：先把汇报边界讲清楚，主文献是 EA 断奶仔猪，Microbiome 只是机制补充。",
            "图怎么看：右侧为文章视觉引入，正文只保留题目、补充机制、汇报人和日期。",
            "结果判断：这里不展开数据，只让老师先知道后面围绕 EA、FEA 和 3-PPA/AhR 三个关键词。",
            "机制解释：主线是植物酚类通过菌群改善氧化应激和屏障，不把两篇文章混成一个实验。",
            "过渡句：下一页先交代汇报路线，让听众知道每一部分解决什么问题。",
        ]
    if t.startswith("汇报路线"):
        return [
            "本页定位：回答这场汇报怎么听，先看主文献证据，再看机制补充，最后落到课题启发。",
            "图怎么看：左右两栏分别代表 Animal Nutrition 主线和 Microbiome 补充线。",
            "结果判断：主文献负责证明 EA 对断奶仔猪有效，补充文献负责解释菌群代谢物如何修复屏障。",
            "机制解释：这避免把高 IF 文章喧宾夺主，也避免把假说讲成已验证结论。",
            "过渡句：路线明确后，先介绍主文献的基本信息和为什么选它主讲。",
        ]
    if t.startswith("文献信息"):
        return [
            "本页定位：回答为什么主讲这篇 Animal Nutrition，而不是直接讲 Microbiome。",
            "图怎么看：看题目、期刊、年份和关键词，重点抓 weanling piglets、ellagic acid、gut microbiota。",
            "结果判断：这篇同时包含断奶仔猪、酚类干预、氧化应激、屏障和菌群，和本人课题贴合度最高。",
            "机制解释：它提供的是课题可借鉴的动物模型和指标体系，而不是单纯背景材料。",
            "过渡句：接下来从断奶应激问题切入，说明为什么需要研究 EA。",
        ]
    if t.startswith("研究背景"):
        return [
            "本页定位：回答断奶仔猪为什么是一个值得关注的肠道健康模型。",
            "图怎么看：按腹泻、生长下降、氧化应激、屏障损伤和菌群紊乱这几条线读。",
            "结果判断：传统抗生素和高锌策略受限后，植物多酚成为潜在替代方向。",
            "机制解释：EA 作为典型植物酚类，可能同时影响抗氧化系统、屏障和菌群。",
            "过渡句：背景之后，下一页把全文科学问题压缩成一个核心假设。",
        ]
    if t.startswith("科学问题"):
        return [
            "本页定位：回答作者真正想证明什么，即 EA 是否通过菌群缓解断奶应激。",
            "图怎么看：先看问题，再看假设，重点记住 EA、菌群、抗氧化和屏障四个节点。",
            "结果判断：如果 EA 有效但 FEA 也能复制部分效果，就说明菌群不是旁观者。",
            "机制解释：这一页建立后面所有 Figure 的判断标准：表型改善是否能被菌群线索解释。",
            "过渡句：下一页看实验设计，尤其是 FEA 组为什么是全文亮点。",
        ]
    if t.startswith("实验设计"):
        return [
            "本页定位：回答三组设计如何支撑菌群介导这个判断。",
            "图怎么看：Ctrl 是基础日粮，EA 是 0.1% EA，FEA 是接受 EA 组粪菌悬液；试验 126 头、14 天。",
            "结果判断：FEA 不是普通对照，而是用来检验 EA 改变后的菌群是否能转移保护效果。",
            "机制解释：只要 FEA 复制腹泻、屏障或氧化应激改善，就能支持菌群介导。",
            "过渡句：设计讲清后，先看最直接的生长性能和腹泻表型。",
        ]
    if t.startswith("Table 3"):
        return [
            "本页定位：回答 EA/FEA 是否先在整体生产表型上有改善迹象。",
            "图怎么看：看 ADG、ADFI、F/G 等生长性能指标，比较 Ctrl、EA、FEA 三组。",
            "结果判断：EA/FEA 对部分生长指标有改善或恢复趋势，说明干预不是只改变分子指标。",
            "机制解释：生长性能是总表型，后面需要用腹泻、屏障和氧化应激解释它为什么改善。",
            "过渡句：下一页进入 Figure 1A，看断奶应激最直观的腹泻评分。",
        ]
    if t.startswith("Figure 1A"):
        return [
            "本页定位：回答 EA 是否真的缓解断奶腹泻。",
            "图怎么看：粪便评分越高代表腹泻越严重，横轴是断奶后第 7 天和第 14 天。",
            "结果判断：EA 组两个时间点评分下降，FEA 组在第 14 天也下降，说明粪菌转移复制了部分效果。",
            "机制解释：这不仅说明 EA 有效，也提示 EA 的作用可能和肠道菌群改变有关。",
            "过渡句：腹泻改善后，下一页看肠屏障蛋白是否同步修复。",
        ]
    if t.startswith("Figure 1B"):
        return [
            "本页定位：回答腹泻改善是否伴随紧密连接屏障修复。",
            "图怎么看：B 是 Western blot 条带，C-D 是 Claudin-1 和 Occludin 定量。",
            "结果判断：EA 和 FEA 主要上调 Occludin，而 Claudin-1 变化不明显。",
            "机制解释：Occludin 上升提示屏障完整性改善，也解释了腹泻评分下降的组织功能基础。",
            "过渡句：屏障改善之外，断奶应激还涉及氧化损伤，下一页看抗氧化能力。",
        ]
    if t.startswith("Figure 1E"):
        return [
            "本页定位：回答 EA 是否提高空肠抗氧化防御能力。",
            "图怎么看：GSH/GSSG 代表还原状态，T-AOC 代表总抗氧化能力，CAT 是抗氧化酶。",
            "结果判断：EA 和 FEA 均提高这些指标，说明抗氧化系统被增强。",
            "机制解释：FEA 的复制效果把抗氧化改善和菌群介导联系起来。",
            "过渡句：抗氧化能力增强后，下一页看氧化损伤产物 MDA 是否下降。",
        ]
    if t.startswith("Figure 1H"):
        return [
            "本页定位：回答 EA 是否真正降低氧化损伤，而不只是提高抗氧化指标。",
            "图怎么看：MDA 是脂质过氧化损伤产物，NO 是氧化/炎症相关指标。",
            "结果判断：EA 和 FEA 降低 MDA，但 NO 变化不显著。",
            "机制解释：结果更支持“降低脂质过氧化损伤”，而不是所有炎症氧化指标都被全面改变。",
            "过渡句：下一页把 Figure 1 的腹泻、屏障和氧化应激证据合在一起。",
        ]
    if t.startswith("Figure 1 小结"):
        return [
            "本页定位：回答 Figure 1 到底建立了哪条初步证据链。",
            "图怎么看：按腹泻评分、Occludin、抗氧化能力和 MDA 四类结果回顾。",
            "结果判断：EA 改善腹泻、上调 Occludin、增强抗氧化并降低 MDA；FEA 复制部分效果。",
            "机制解释：这支持 EA 的表型保护和菌群介导可能性，但还需要组织形态证据。",
            "过渡句：所以下一部分进入 Figure 2，看空肠结构和凋亡。",
        ]
    if t.startswith("Figure 2："):
        return [
            "本页定位：把问题从功能指标推进到组织学层面。",
            "图怎么看：后面会依次看 HE、TUNEL、绒毛高度、隐窝深度、绒隐比和凋亡定量。",
            "结果判断：如果 EA 改善这些指标，就说明保护作用不只是生化变化，也有形态学支撑。",
            "机制解释：形态修复与屏障修复共同解释断奶腹泻缓解。",
            "过渡句：先看 Fig.2A 的 HE 和 TUNEL 原始图像。",
        ]
    if t.startswith("Figure 2A"):
        return [
            "本页定位：回答 EA 是否改善空肠组织形态并减少细胞凋亡。",
            "图怎么看：上排 HE 看绒毛结构，中下排 TUNEL/DAPI 看绿色凋亡信号。",
            "结果判断：Ctrl 组结构受损、TUNEL 信号较多，EA 和 FEA 组绒毛更完整、凋亡减少。",
            "机制解释：这说明 EA 的屏障保护有组织学基础，FEA 仍提示菌群参与。",
            "过渡句：下一页用 B-D 的定量指标验证形态改善。",
        ]
    if t.startswith("Figure 2B"):
        return [
            "本页定位：回答 EA 是否改善空肠吸收结构。",
            "图怎么看：B 是绒毛高度，C 是隐窝深度，D 是绒毛高度/隐窝深度比值。",
            "结果判断：EA 提高绒毛高度和绒隐比，隐窝深度变化不明显。",
            "机制解释：绒毛结构改善意味着吸收面积和屏障状态恢复，是腹泻缓解的形态支撑。",
            "过渡句：形态改善之后，下一页单独看 TUNEL AOD 的凋亡定量。",
        ]
    if t.startswith("Figure 2E"):
        return [
            "本页定位：回答 EA 是否降低上皮细胞凋亡。",
            "图怎么看：TUNEL staining AOD 越高代表凋亡信号越强，三组并列比较。",
            "结果判断：EA 降低 TUNEL AOD，FEA 也有下降趋势或部分复制效果。",
            "机制解释：凋亡减少说明 EA 可减轻断奶应激导致的上皮损伤，进一步支持屏障修复。",
            "过渡句：下一页把 Figure 2 的形态和凋亡证据合并总结。",
        ]
    if t.startswith("Figure 2 小结"):
        return [
            "本页定位：回答 Figure 2 对主线贡献是什么。",
            "图怎么看：把 HE/TUNEL 图像和 B-E 定量结果放在一起读。",
            "结果判断：EA 改善绒毛结构、提高绒隐比并降低凋亡，FEA 部分复制。",
            "机制解释：这说明 EA 的保护不是停留在粪便评分，而是落实到上皮结构修复。",
            "过渡句：接下来进入 Figure 3，看看抗氧化和炎症基因能否解释这种修复。",
        ]
    if t.startswith("Figure 3："):
        return [
            "本页定位：把组织学改善进一步追问到转录层面。",
            "图怎么看：后面分两组讲，先看抗氧化基因，再看炎症因子。",
            "结果判断：本 Figure 的判断重点是哪个通路更突出，而不是每个因子都显著。",
            "机制解释：如果抗氧化基因变化强于炎症因子，就说明 EA 主机制更偏抗氧化。",
            "过渡句：先看 HO-1 和 NQO-1。",
        ]
    if t.startswith("Figure 3A"):
        return [
            "本页定位：回答 EA 是否激活经典抗氧化相关基因。",
            "图怎么看：A 是 HO-1，B 是 NQO-1，纵轴为相对 mRNA 表达。",
            "结果判断：EA 明显上调 HO-1，NQO-1 变化不明显。",
            "机制解释：抗氧化转录响应具有选择性，不能简单说所有 Nrf2 下游都全面激活。",
            "过渡句：下一页看与谷胱甘肽合成更直接相关的 GCLC 和 GCLM。",
        ]
    if t.startswith("Figure 3C"):
        return [
            "本页定位：回答 EA 是否增强谷胱甘肽合成相关防御。",
            "图怎么看：GCLC 和 GCLM 是谷胱甘肽合成相关关键基因，与 Fig.1 的 GSH/GSSG 对应。",
            "结果判断：EA 上调 GCLC/GCLM，FEA 对 GCLM 也有提升。",
            "机制解释：这把抗氧化酶活和转录层面的谷胱甘肽系统连接起来。",
            "过渡句：抗氧化证据清楚后，下一页看炎症因子是否同样强。",
        ]
    if t.startswith("Figure 3E"):
        return [
            "本页定位：回答 EA 是否主要通过抑制促炎因子发挥作用。",
            "图怎么看：TNF-α、IL-6、IL-1α 是促炎相关 mRNA 指标。",
            "结果判断：多数促炎因子没有显著下降，说明炎症证据弱于抗氧化证据。",
            "机制解释：汇报时要主动讲这个局限，避免把 EA 过度解释成强抗炎。",
            "过渡句：下一页补充 IL-1β 和 IL-10，看是否有温和免疫调节。",
        ]
    if t.startswith("Figure 3H"):
        return [
            "本页定位：补充判断炎症通路是否存在方向性变化。",
            "图怎么看：IL-1β 属促炎，IL-10 属抗炎，仍然看三组相对表达。",
            "结果判断：IL-1β 变化不明显，IL-10 在 FEA 中有上升趋势。",
            "机制解释：这提示菌群可能有温和免疫调节，但不是本文最强证据。",
            "过渡句：下一页总结 Figure 3，把抗氧化作为主线讲清。",
        ]
    if t.startswith("Figure 3 小结"):
        return [
            "本页定位：回答 Figure 3 应该怎么被老师理解。",
            "图怎么看：抗氧化基因和炎症因子分开看，避免把所有基因混成一个结论。",
            "结果判断：HO-1、GCLC、GCLM 证据更强，炎症因子变化有限。",
            "机制解释：EA 主要支持增强抗氧化防御，而不是强烈抑炎。",
            "过渡句：由于 FEA 多次复制效果，下一部分必须进入菌群结构。",
        ]
    if t.startswith("Figure 4："):
        return [
            "本页定位：正式回答 EA 是否重塑肠道菌群。",
            "图怎么看：后面先看测序深度和 OTU 共享，再看 PCoA 与 Alpha diversity。",
            "结果判断：如果 EA 和 FEA 结构接近，就支持菌群在保护作用中的中介地位。",
            "机制解释：菌群变化是从表型保护走向机制解释的关键桥梁。",
            "过渡句：先确认测序数据可靠，再谈菌群差异。",
        ]
    if t.startswith("Figure 4A"):
        return [
            "本页定位：回答 16S 数据是否可靠，以及 EA/FEA 是否共享更多菌群特征。",
            "图怎么看：A 是稀释曲线，趋平说明测序深度足够；B 是 OTU Venn 图。",
            "结果判断：EA 与 FEA 共享 OTUs 较多，说明 FEA 菌群向 EA 组靠近。",
            "机制解释：这为“EA 改变菌群，菌群可转移保护效应”提供结构证据。",
            "过渡句：下一页看 PCoA，判断整体菌群结构是否真正分离。",
        ]
    if t.startswith("Figure 4C"):
        return [
            "本页定位：回答 EA/FEA 是否显著改变整体菌群结构。",
            "图怎么看：PCoA 看 beta diversity 点群分离，Table 4 看 alpha diversity 指标。",
            "结果判断：三组菌群结构显著分离，EA 提高 Observed species，其他指标按显著性谨慎解释。",
            "机制解释：这说明 EA 不只影响单个菌，而是重塑群落结构。",
            "过渡句：下一页总结 Figure 4 后，再看具体差异菌是谁。",
        ]
    if t.startswith("Figure 4 小结"):
        return [
            "本页定位：回答 Figure 4 对菌群介导假说的贡献。",
            "图怎么看：依次回顾稀释曲线、OTU 共享、PCoA 和 alpha diversity。",
            "结果判断：EA/FEA 改变菌群结构，FEA 与 EA 趋势相近。",
            "机制解释：这支持菌群是 EA 保护作用的重要中介，但还需要具体菌和代谢物证据。",
            "过渡句：下一部分进入 Figure 5，寻找关键差异菌。",
        ]
    if t.startswith("Figure 5："):
        return [
            "本页定位：把“菌群改变”进一步拆到哪些菌改变。",
            "图怎么看：后面从门水平、种水平和作者筛出的差异菌逐步讲。",
            "结果判断：不需要背所有菌名，重点看哪些菌与 SCFAs、屏障和氧化应激可能相关。",
            "机制解释：关键菌是连接 16S 结构变化和代谢功能变化的中间层。",
            "过渡句：先看门水平和种水平的整体组成。",
        ]
    if t.startswith("Figure 5A"):
        return [
            "本页定位：回答 EA 是否改变菌群组成格局。",
            "图怎么看：A 是门水平组成，B 是种水平组成，堆叠柱图看相对丰度比例。",
            "结果判断：EA/FEA 使主要菌门和物种比例发生偏移。",
            "机制解释：这说明 EA 对菌群不是微弱扰动，而是有群落组成层面的影响。",
            "过渡句：下一页从组成图转到显著差异菌。",
        ]
    if t.startswith("Figure 5C"):
        return [
            "本页定位：回答哪些分类单元首先表现出显著变化。",
            "图怎么看：C-D 分别看 Tenericutes 和 Parabacteroides 的相对丰度。",
            "结果判断：EA/FEA 对这两类菌有方向性调节，说明具体菌群响应不同。",
            "机制解释：这些差异菌后面可与宿主指标做相关性，形成机制线索。",
            "过渡句：下一页重点看与发酵和 SCFAs 更相关的 Ruminococcaceae。",
        ]
    if t.startswith("Figure 5E"):
        return [
            "本页定位：回答 EA 是否富集可能与发酵代谢相关的菌群。",
            "图怎么看：E 看 Ruminococcaceae，F 看 Enterobacteriaceae，仍以三组相对丰度比较。",
            "结果判断：Ruminococcaceae 上升，FEA 组也出现相关变化；Enterobacteriaceae 需谨慎解释。",
            "机制解释：Ruminococcaceae 为后面 SCFAs 增加提供菌群来源线索。",
            "过渡句：下一页继续看种水平差异菌和相关性线索。",
        ]
    if t.startswith("Figure 5G"):
        return [
            "本页定位：回答 Clostridium ramosum 等差异菌是否值得关注。",
            "图怎么看：G-H 比较 Streptococcaceae 与 C. ramosum 的相对丰度。",
            "结果判断：这些菌在 EA/FEA 间出现显著差异，是 Figure 7 相关性分析的候选对象。",
            "机制解释：差异菌本身不是结论，关键是它们能否连接抗氧化、屏障和代谢物。",
            "过渡句：下一页补充 FEA 组中特别明显的种水平变化。",
        ]
    if t.startswith("Figure 5I"):
        return [
            "本页定位：回答 FEA 是否带来与屏障相关的种水平菌变化。",
            "图怎么看：I-J 分别看 Veillonella parvula 和 Akkermansia muciniphila。",
            "结果判断：FEA 组中两类菌增加，提示粪菌移植确实重塑了菌群。",
            "机制解释：Akkermansia 常与黏液层和屏障有关，但汇报时仍以原文数据为准。",
            "过渡句：下一页总结 Figure 5，再进入菌群功能和代谢物。",
        ]
    if t.startswith("Figure 5 小结"):
        return [
            "本页定位：回答差异菌结果如何服务主线。",
            "图怎么看：从整体组成到具体菌群，关注能和 SCFAs、屏障、氧化应激连接的菌。",
            "结果判断：EA/FEA 调节多类关键菌，FEA 与 EA 的相似性支持菌群介导。",
            "机制解释：菌群重塑需要进一步看功能代谢物，否则仍停留在相关层面。",
            "过渡句：下一部分进入 Figure 6 和 Table 5，看功能预测与 SCFAs。",
        ]
    if t.startswith("Figure 6："):
        return [
            "本页定位：回答菌群结构改变是否转化成功能代谢改变。",
            "图怎么看：后面先看 KEGG 功能预测，再看短链脂肪酸实测结果。",
            "结果判断：如果功能通路和 SCFAs 同时变化，菌群机制会更完整。",
            "机制解释：代谢物是菌群影响宿主屏障和抗氧化状态的重要媒介。",
            "过渡句：先看 EA 相对 Ctrl 的功能预测差异。",
        ]
    if t.startswith("Figure 6A"):
        return [
            "本页定位：回答 EA 是否改变菌群预测功能通路。",
            "图怎么看：横向条形图比较 Ctrl 与 EA 的 KEGG 功能通路丰度和显著性。",
            "结果判断：EA 改变多类代谢相关通路，说明菌群变化可能具有功能后果。",
            "机制解释：这为 SCFAs 和其他代谢物变化提供功能层面的铺垫。",
            "过渡句：下一页看 FEA 是否复制类似功能变化。",
        ]
    if t.startswith("Figure 6B"):
        return [
            "本页定位：回答粪菌移植是否复制 EA 的菌群功能变化。",
            "图怎么看：比较 Ctrl 与 FEA 的预测通路差异，关注与 EA 结果是否同向。",
            "结果判断：FEA 复制部分功能变化，但不等同于 EA 全部效应。",
            "机制解释：这说明菌群介导是重要部分，但 EA 也可能有直接宿主作用。",
            "过渡句：下一页看实测 SCFAs，把功能预测落到代谢物层面。",
        ]
    if t.startswith("Table 5"):
        return [
            "本页定位：回答菌群变化是否产生可检测的功能代谢物改变。",
            "图怎么看：Table 5 比较空肠和结肠内容物中乙酸、丙酸、丁酸等 SCFAs。",
            "结果判断：EA/FEA 提高部分 SCFAs，尤其结肠内容物变化更明显。",
            "机制解释：SCFAs 可连接菌群发酵、屏障保护和免疫/抗氧化状态。",
            "过渡句：下一页总结 Figure 6，随后看相关性是否支持这条链。",
        ]
    if t.startswith("Figure 6 小结"):
        return [
            "本页定位：回答菌群功能和 SCFAs 结果如何补强机制。",
            "图怎么看：把功能预测和 Table 5 的代谢物检测一起理解。",
            "结果判断：EA/FEA 不仅改变菌群组成，还改变部分代谢功能和 SCFAs。",
            "机制解释：这支持“菌群—代谢物—屏障/抗氧化”的中间环节。",
            "过渡句：最后用 Figure 7 看差异菌与宿主表型是否相关。",
        ]
    if t.startswith("Figure 7"):
        return [
            "本页定位：回答差异菌是否和宿主表型、SCFAs、屏障和氧化应激指标有关。",
            "图怎么看：热图横轴是宿主指标，纵轴是差异菌，颜色表示相关方向，星号表示显著性。",
            "结果判断：多种差异菌与 MDA、TUNEL、Occludin、GSH/GSSG 和 SCFAs 显著相关。",
            "机制解释：这是相关性证据，支持菌群介导，但不能等同于因果证明。",
            "过渡句：下一页总结主文献，把强证据和局限都讲出来。",
        ]
    if t.startswith("主文献总结"):
        return [
            "本页定位：用一页讲清 Animal Nutrition 的核心证据链。",
            "图怎么看：按腹泻下降、屏障改善、氧化应激下降、菌群和 SCFAs 改变的顺序回顾。",
            "结果判断：EA 有效，FEA 复制部分效果，说明菌群参与是全文最重要的设计亮点。",
            "机制解释：主文献支持 EA—菌群—抗氧化—屏障，但没有直接证明 3-PPA/AhR。",
            "过渡句：正因为这个空白，下一部分引入 Microbiome 作为机制补充。",
        ]
    if t.startswith("机制补充"):
        return [
            "本页定位：回答为什么主文献后面还要讲另一篇文章。",
            "图怎么看：这是一页过渡，不展开新数据，只指出主文献没有讲透的机制问题。",
            "结果判断：Animal Nutrition 有菌群和 SCFAs，但没有直接找到芳香族代谢物和宿主受体通路。",
            "机制解释：Microbiome 提供 B. fragilis—3-PPA—AhR—屏障这条机制链。",
            "过渡句：下一页先介绍 Microbiome 文章本身，避免两篇文章混淆。",
        ]
    if t.startswith("Microbiome 文献信息"):
        return [
            "本页定位：说明 Microbiome 是独立机制补充文献，不是 EA 研究。",
            "图怎么看：看题目中的 pig gut microbiota、3-PPA、AhR 和 intestinal barrier。",
            "结果判断：它研究的是猪源菌群代谢物修复屏障，不是酚类干预，也不是断奶仔猪模型。",
            "机制解释：因此只能用来补充“菌群代谢物如何作用于宿主”，不能替代主文献结论。",
            "过渡句：下一页明确它补的机制空白是什么。",
        ]
    if t.startswith("为什么引入 Microbiome"):
        return [
            "本页定位：回答引入高 IF 文献的必要性。",
            "图怎么看：只看三个问题：哪个菌、哪个代谢物、哪条宿主信号通路。",
            "结果判断：Animal Nutrition 证明 EA 调菌群，但 Microbiome 进一步证明 3-PPA 可通过 AhR 修复屏障。",
            "机制解释：这让汇报从相关性走向可验证的代谢物—受体机制假说。",
            "过渡句：下一页先看 Microbiome 如何证明菌群能转移屏障表型。",
        ]
    if t.startswith("Microbiome Fig.2a"):
        return [
            "本页定位：回答不同猪源菌群是否与屏障功能差异有关。",
            "图怎么看：a-c 看 DAO、内毒素、D-乳酸等屏障损伤指标，d 展示无菌小鼠 FMT 设计。",
            "结果判断：CM 猪屏障相关指标更好，作者进一步用 FMT 排除单纯宿主背景影响。",
            "机制解释：这说明菌群本身可能携带屏障保护能力。",
            "过渡句：下一页看移植后损伤指标是否真的改善。",
        ]
    if t.startswith("Microbiome Fig.2e"):
        return [
            "本页定位：回答 CM 菌群移植是否降低屏障损伤指标。",
            "图怎么看：e-g 比较 DLY-R 与 CM-R 小鼠的 DAO、内毒素和 D-乳酸。",
            "结果判断：CM-R 组这些损伤或通透性指标下降，说明菌群表型可以被转移。",
            "机制解释：这为“菌群可以调控屏障功能”提供功能层面的因果线索。",
            "过渡句：下一页看屏障结构蛋白是否同步提高。",
        ]
    if t.startswith("Microbiome Fig.2h"):
        return [
            "本页定位：回答 CM 菌群是否增强上皮连接结构。",
            "图怎么看：h 是 Western blot，i-k 是 ZO-1、E-cadherin、Connexin 43 定量。",
            "结果判断：CM-R 组连接蛋白升高，和损伤指标下降方向一致。",
            "机制解释：这说明菌群改善屏障不是只改变血清指标，而是落实到连接蛋白。",
            "过渡句：下一页补充黏膜免疫指标，形成更完整的屏障表型。",
        ]
    if t.startswith("Microbiome Fig.2l"):
        return [
            "本页定位：补充说明 CM 菌群还影响黏膜免疫或局部防御状态。",
            "图怎么看：看 IgA、黏膜相关指标或免疫指标在 DLY-R 与 CM-R 间的差异。",
            "结果判断：CM 菌群移植后部分防御指标改善，与连接蛋白结果一致。",
            "机制解释：菌群对屏障的影响包括结构连接和黏膜防御两个层面。",
            "过渡句：接下来追问是哪种菌和代谢物承担这种作用。",
        ]
    if t.startswith("Microbiome Fig.7a"):
        return [
            "本页定位：回答候选代谢物中谁最可能直接改善屏障。",
            "图怎么看：作者用多种代谢物处理后，比较屏障损伤指标变化。",
            "结果判断：3-PPA 在候选物中表现突出，可降低屏障损伤相关指标。",
            "机制解释：这把菌群作用从“菌名变化”推进到“代谢物功能”。",
            "过渡句：下一页看 3-PPA 是否真正提高连接蛋白。",
        ]
    if t.startswith("Microbiome Fig.7h"):
        return [
            "本页定位：回答 3-PPA 是否能直接增强屏障结构。",
            "图怎么看：h 是连接蛋白 Western blot，i-k 是 ZO-1、E-cadherin、Connexin 43 定量。",
            "结果判断：3-PPA 提高多种连接蛋白，说明它不是单纯相关代谢物。",
            "机制解释：这为菌群代谢物修复屏障提供直接功能证据。",
            "过渡句：下一页看 3-PPA 与 B. fragilis 的来源关系。",
        ]
    if t.startswith("Microbiome Fig.7l"):
        return [
            "本页定位：回答 3-PPA 可能来自哪类关键菌。",
            "图怎么看：看代谢组差异、菌株处理和 3-PPA 水平变化。",
            "结果判断：B. fragilis 与 3-PPA 富集形成关键菌—代谢物线索。",
            "机制解释：这把机制链前半段从“菌群”具体化为 B. fragilis 和 3-PPA。",
            "过渡句：下一页进入 AhR，看 3-PPA 通过哪条宿主通路起效。",
        ]
    if t.startswith("Microbiome Fig.8a"):
        return [
            "本页定位：回答 3-PPA 或 B. fragilis 是否激活 AhR 信号。",
            "图怎么看：看 AhR 核转位、CYP1A1 等 AhR 下游响应指标。",
            "结果判断：3-PPA 和 B. fragilis 均增强 AhR 相关信号。",
            "机制解释：这把代谢物功能连接到宿主受体通路，为屏障修复提供机制解释。",
            "过渡句：下一页看抑制 AhR 后保护作用是否消失。",
        ]
    if t.startswith("Microbiome Fig.8i"):
        return [
            "本页定位：回答 AhR 是否是 3-PPA 保护屏障的必要通路。",
            "图怎么看：加入 AhR 抑制剂后再看屏障损伤指标和通路指标。",
            "结果判断：抑制 AhR 后，3-PPA 的保护作用被削弱。",
            "机制解释：这是比相关性更强的阻断证据，说明 AhR 位于机制链关键位置。",
            "过渡句：下一页再看连接蛋白是否也被阻断。",
        ]
    if t.startswith("Microbiome Fig.8p"):
        return [
            "本页定位：回答 AhR 被抑制后，连接蛋白改善是否还能维持。",
            "图怎么看：比较 3-PPA 单独处理和 3-PPA 加 AhR 抑制剂后的连接蛋白水平。",
            "结果判断：AhR 抑制削弱 ZO-1、E-cadherin、Connexin 43 的改善。",
            "机制解释：这完整支持 B. fragilis—3-PPA—AhR—屏障这条链。",
            "过渡句：下一页把两篇文献整合为本人课题假说。",
        ]
    if t.startswith("整合机制假说"):
        return [
            "本页定位：把两篇文献连接成一个可验证假说，同时避免过度下结论。",
            "图怎么看：分三层读：主文献已验证、Microbiome 机制补充、本人课题待验证。",
            "结果判断：EA 改善菌群和屏障是主文献证据，3-PPA/AhR 是独立文献证据。",
            "机制解释：荔枝酚类是否走芳香族代谢物/SCFAs—AhR/Nrf2—屏障修复，需要后续实验验证。",
            "过渡句：下一页把这个假说转成课题指标设计。",
        ]
    if t.startswith("课题启发"):
        return [
            "本页定位：回答这两篇文献对荔枝酚类课题具体有什么用。",
            "图怎么看：按动物模型、干预物、表型指标、菌群代谢物和信号通路逐项看。",
            "结果判断：可以直接借鉴断奶仔猪模型和 EA 的指标体系，同时补充芳香族代谢物与 AhR/Nrf2。",
            "机制解释：重点从“现象改善”升级到“酚类—菌群—代谢物—宿主通路”的可验证链条。",
            "过渡句：下一页用几句话收束整场汇报的亮点。",
        ]
    if t.startswith("研究总结"):
        return [
            "本页定位：用老师能记住的方式总结本次汇报。",
            "图怎么看：分主文献结论、机制补充和汇报亮点三块读。",
            "结果判断：主文献证据链完整，FEA 组是亮点；Microbiome 提供更深的代谢物—AhR 机制。",
            "机制解释：两篇文章共同服务于荔枝酚类课题假说，但不能混为同一实验证据。",
            "过渡句：后面附上两天阅读路线和答辩策略，帮助准备老师提问。",
        ]
    if t.startswith("两天内读懂"):
        return [
            "本页定位：给自己准备汇报的阅读路线，先把 Animal Nutrition 主文献吃透。",
            "图怎么看：左侧按 Abstract、Methods、Figure、Discussion 顺序读，右侧写 Day 1 的 3 分钟目标。",
            "结果判断：第一天必须讲清 Ctrl、EA、FEA、126 头、0.1% EA 和 14 天设计。",
            "机制解释：Day 1 的重点是建立 EA—菌群—抗氧化—屏障主证据链。",
            "过渡句：下一页只读 Microbiome 中与 3-PPA/AhR 有关的机制图。",
        ]
    if t.startswith("Day 2"):
        return [
            "本页定位：给第二天的阅读任务划边界，避免把 Microbiome 展开成第二篇主讲。",
            "图怎么看：只抓 pig gut microbiota、B. fragilis、3-PPA、AhR 和 intestinal barrier。",
            "结果判断：需要能用 2 分钟解释为什么补充这篇，以及它怎样提高机制深度。",
            "机制解释：Microbiome 的作用是提供菌群代谢物—AhR—屏障这条可借鉴机制。",
            "过渡句：下一页总结面向博士和老师时应该怎么讲、哪些坑要主动避开。",
        ]
    if t.startswith("面向博士"):
        return [
            "本页定位：提醒自己在老师面前要讲主线、讲边界、讲局限。",
            "图怎么看：按五条策略读，尤其注意不要把两篇文章说成同一个研究。",
            "结果判断：主文献支持 EA 通过菌群改善屏障和氧化应激，荔枝酚类和 3-PPA/AhR 属于待验证假说。",
            "机制解释：主动讲局限会显得更像研究生读文献，而不是只报好结果。",
            "过渡句：下一页准备几个老师最可能追问的问题。",
        ]
    if t.startswith("可能被老师"):
        return [
            "本页定位：提前准备答问，避免现场把主文献和机制补充讲混。",
            "图怎么看：五个问题分别对应选题理由、FEA 意义、3-PPA/AhR 边界、课题借鉴和两文关系。",
            "结果判断：回答时始终强调 Animal Nutrition 是主线，Microbiome 是独立机制补充。",
            "机制解释：所有延伸到荔枝酚类的内容都要说成待验证假说，而不是已有结论。",
            "过渡句：最后回到致谢页，正式结束汇报并邀请老师提建议。",
        ]
    if t.startswith("致"):
        return [
            "本页定位：结束汇报，给老师提问留出空间。",
            "图怎么看：页面只保留感谢语，不再加入新信息，避免结尾分散注意力。",
            "结果判断：如果老师追问，优先回到五句话主线：EA 有效、FEA 证明菌群参与、Microbiome 补机制、荔枝酚类待验证。",
            "机制解释：答问时坚持区分直接证据、机制补充和个人课题假说。",
            "过渡句：可以说感谢各位老师批评指正，欢迎老师对后续实验设计提出建议。",
        ]
    return [
        f"本页定位：围绕“{t}”回答本模块中的一个具体问题。",
        "图怎么看：先说明分组、横纵坐标和检测指标含义，再指出显著性标记。",
        "结果判断：比较 EA、FEA 与 Ctrl 的方向变化，注意哪些结果是部分复制。",
        "机制解释：把结果放回抗氧化、屏障修复、菌群介导或代谢物调节这条链中。",
        "过渡句：下一页继续沿着证据链向更深一层机制推进。",
    ]


def patch_notes(pptx_path: Path, out_path: Path, titles: list[str]):
    ns_rel = "http://schemas.openxmlformats.org/package/2006/relationships"
    rel_tag = f"{{{ns_rel}}}Relationship"
    with zipfile.ZipFile(pptx_path, "r") as zin:
        files = {name: zin.read(name) for name in zin.namelist()}

    pres = ET.fromstring(files["ppt/presentation.xml"])
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main", "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}
    rels = ET.fromstring(files["ppt/_rels/presentation.xml.rels"])
    rid_to_target = {el.attrib["Id"]: el.attrib["Target"] for el in rels}
    order = []
    for sld in pres.find("p:sldIdLst", ns):
        rid = sld.attrib[f"{{{ns['r']}}}id"]
        order.append("ppt/" + rid_to_target[rid])

    existing_note_nums = []
    for name in files:
        m = re.match(r"ppt/notesSlides/notesSlide(\d+)\.xml$", name)
        if m:
            existing_note_nums.append(int(m.group(1)))
    next_note = max(existing_note_nums or [0]) + 1

    content_types = files["[Content_Types].xml"].decode("utf-8")

    for idx, slide_part in enumerate(order, 1):
        rel_path = "ppt/slides/_rels/" + Path(slide_part).name + ".rels"
        rel_xml = files.get(rel_path, b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>').decode("utf-8")
        rel_root = ET.fromstring(rel_xml)
        note_target = None
        max_rid = 0
        for el in rel_root.findall(rel_tag):
            rid = el.attrib.get("Id", "")
            if rid.startswith("rId"):
                try:
                    max_rid = max(max_rid, int(rid[3:]))
                except ValueError:
                    pass
            if el.attrib.get("Type", "").endswith("/notesSlide"):
                note_target = el.attrib["Target"]
        if note_target is None:
            note_name = f"notesSlide{next_note}.xml"
            next_note += 1
            note_target = f"../notesSlides/{note_name}"
            ET.SubElement(
                rel_root,
                rel_tag,
                {
                    "Id": f"rId{max_rid + 1}",
                    "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
                    "Target": note_target,
                },
            )
            files[rel_path] = ET.tostring(rel_root, encoding="utf-8", xml_declaration=True)
        note_part = "ppt/notesSlides/" + Path(note_target).name
        note_rels = "ppt/notesSlides/_rels/" + Path(note_target).name + ".rels"
        lines = note_lines(titles[idx - 1])
        if note_part in files:
            files[note_part] = patch_note_text_xml(files[note_part].decode("utf-8"), lines, idx).encode("utf-8")
        else:
            files[note_part] = make_note_xml(lines, idx).encode("utf-8")
            part_name = "/" + note_part
            if part_name not in content_types:
                content_types = content_types.replace(
                    "</Types>",
                    f'<Override PartName="{part_name}" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/></Types>',
                )
        if note_rels not in files:
            slide_target = "../slides/" + Path(slide_part).name
            files[note_rels] = (
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="../notesMasters/notesMaster1.xml"/>'
                f'<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="{slide_target}"/>'
                "</Relationships>"
            ).encode("utf-8")

    files["[Content_Types].xml"] = content_types.encode("utf-8")
    with zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in files.items():
            zout.writestr(name, data)


def main():
    ASSETS.mkdir(parents=True, exist_ok=True)
    bd_img, e_img = crop_fig2_panels()

    prs = Presentation(SRC)

    title_updates = {
        "汇报整体逻辑：主文献讲证据，半篇文献讲机制补充": "汇报路线：先证明 EA 有效，再解释菌群代谢物机制",
        "主讲文献信息：Animal Nutrition 鞣花酸-断奶仔猪文章": "文献信息：Animal Nutrition 主文献聚焦 EA 与断奶仔猪肠道健康",
        "研究背景：断奶应激导致肠道损伤，抗生素替代需求增加": "研究背景：断奶应激导致腹泻、氧化应激和屏障损伤",
        "核心科学问题与研究假设": "科学问题：EA 是否通过肠道菌群缓解断奶应激？",
        "实验设计：Ctrl、EA、FEA 三组，14 天干预": "实验设计：Ctrl、EA、FEA 三组用于判断菌群介导作用",
        "Table 3：EA/FEA 对生长性能的影响": "Table 3：EA/FEA 改善部分生长性能，提示整体保护作用",
        "Figure 1A：EA/FEA 降低断奶仔猪粪便评分": "Figure 1A：EA 降低粪便评分，提示腹泻缓解",
        "Figure 1B-D：EA/FEA 改善紧密连接蛋白 Occludin": "Figure 1B-D：EA 和 FEA 上调 Occludin，提示肠屏障改善",
        "Figure 1E-G：EA/FEA 提高空肠抗氧化能力": "Figure 1E-G：EA 增强空肠抗氧化能力",
        "Figure 1H-I：EA/FEA 降低 MDA，NO 无显著变化": "Figure 1H-I：EA 降低 MDA，提示脂质过氧化损伤减轻",
        "Figure 1 小结：EA 有效，FEA 支持菌群参与": "Figure 1 小结：EA 初步改善腹泻、屏障和氧化应激",
        "Figure 2：EA/FEA 改善肠道形态与上皮凋亡": "Figure 2：EA 是否进一步修复空肠形态和上皮损伤？",
        "Figure 2A：HE 与 TUNEL 显示空肠结构和凋亡改善": "Figure 2A：EA 改善空肠组织形态并减少 TUNEL 阳性信号",
        "Figure 2B-E：EA 提高绒毛高度和绒隐比，降低凋亡": "Figure 2B-D：EA 提高绒毛高度和绒隐比，改善吸收结构",
        "Figure 2 小结：组织学支持 EA 修复肠道损伤": "Figure 2 小结：EA 修复空肠形态，FEA 支持菌群参与",
        "Figure 3：EA/FEA 调控抗氧化与炎症相关基因": "Figure 3：抗氧化基因是 EA 作用的主要转录证据",
        "Figure 3A-B：HO-1 上调，NQO-1 变化不明显": "Figure 3A-B：EA 上调 HO-1，NQO-1 变化不明显",
        "Figure 3C-D：GCLC/GCLM 上调，提示谷胱甘肽合成增强": "Figure 3C-D：EA 上调 GCLC/GCLM，提示谷胱甘肽合成增强",
        "Figure 3E-G：促炎因子 TNF-α、IL-6、IL-1α 变化不明显": "Figure 3E-G：促炎因子变化有限，炎症不是最强主线",
        "Figure 3H-I：IL-10 有上升趋势，但整体炎症证据较弱": "Figure 3H-I：IL-10 有上升趋势，提示温和免疫调节",
        "Figure 3 小结：EA 的转录层面证据更偏抗氧化而非强抗炎": "Figure 3 小结：EA 更偏向增强抗氧化而非强抑炎",
        "Figure 4：EA/FEA 改变断奶仔猪肠道菌群多样性": "Figure 4：EA 是否重塑断奶仔猪肠道菌群？",
        "Figure 4A-B：测序深度足够，EA 与 FEA 共享更多 OTUs": "Figure 4A-B：测序深度充分，EA 与 FEA 共享更多 OTUs",
        "Figure 4C + Table 4：PCoA 显示菌群结构分离，Alpha 多样性部分提升": "Figure 4C/Table 4：EA 和 FEA 显著改变菌群结构",
        "Figure 4 小结：EA 改变菌群，FEA 组支持菌群参与保护作用": "Figure 4 小结：EA/FEA 重塑菌群，支持菌群介导保护",
        "Figure 5：EA/FEA 改变关键菌群丰度": "Figure 5：哪些关键菌群响应 EA 和 FEA？",
        "Figure 5A-B：门水平与种水平菌群组成改变": "Figure 5A-B：EA 改变门水平和种水平菌群组成",
        "Figure 5C-D：Tenericutes 增加，Parabacteroides 下降": "Figure 5C-D：EA 调节 Tenericutes 和 Parabacteroides",
        "Figure 5E-F：Ruminococcaceae 上升，Enterobacteriaceae 在 FEA 组升高": "Figure 5E-F：Ruminococcaceae 上升，连接 SCFAs 线索",
        "Figure 5G-H：Streptococcaceae 与 Clostridium ramosum 改变": "Figure 5G-H：Clostridium ramosum 等差异菌提示代谢关联",
        "Figure 5I-J：Veillonella parvula 与 Akkermansia muciniphila 在 FEA 组增加": "Figure 5I-J：FEA 增加 Veillonella parvula 和 Akkermansia",
        "Figure 5 小结：关键菌变化连接到 SCFAs 和宿主指标": "Figure 5 小结：差异菌为 SCFAs 和表型相关性提供线索",
        "Figure 6：EA/FEA 调节菌群预测功能与短链脂肪酸": "Figure 6：菌群变化是否转化为功能代谢改变？",
        "Figure 6A：EA 改变菌群功能预测通路": "Figure 6A：EA 改变菌群预测功能通路",
        "Figure 6B：FEA 复制部分菌群功能变化": "Figure 6B：FEA 复制部分功能变化，支持菌群中介",
        "Table 5：EA/FEA 提高结肠和空肠 SCFAs": "Table 5：EA/FEA 提高部分 SCFAs，连接菌群和屏障",
        "Figure 6 小结：菌群变化进一步反映到代谢功能和 SCFAs": "Figure 6 小结：SCFAs 增加支持菌群代谢物参与保护",
        "Figure 7：差异菌与宿主指标的相关性分析": "Figure 7：差异菌与抗氧化、屏障和 SCFAs 指标相关",
        "主文献总结：EA 通过菌群改善断奶仔猪肠道健康": "主文献总结：EA—菌群—抗氧化—屏障构成核心证据链",
        "机制补充：Microbiome 3-PPA/AhR 猪肠屏障研究": "机制补充：主文献未回答菌群代谢物如何修复屏障",
        "Microbiome 文献信息：猪源菌群代谢物 3-PPA 通过 AhR 改善屏障": "Microbiome 文献信息：3-PPA 通过 AhR 增强猪肠屏障",
        "为什么引入 Microbiome：补上主文献的机制空白": "为什么引入 Microbiome：补上“菌群代谢物—宿主通路”",
        "Microbiome Fig.2a-d：CM 猪屏障更强，并设计粪菌移植验证": "Microbiome Fig.2a-d：CM 猪屏障更强，FMT 用于验证菌群作用",
        "Microbiome Fig.2e-k：CM 菌群可转移屏障优势并提高连接蛋白": "Microbiome Fig.2e-g：CM 菌群移植降低屏障损伤指标",
        "Microbiome Fig.2l-r：CM 菌群还影响黏膜免疫相关指标": "Microbiome Fig.2l-r：CM 菌群进一步影响黏膜免疫指标",
        "Microbiome Fig.7a-g：筛选关键代谢物 3-PPA": "Microbiome Fig.7a-g：筛选出屏障保护相关代谢物 3-PPA",
        "Microbiome Fig.7h-k：3-PPA 提高肠屏障连接蛋白": "Microbiome Fig.7h-k：3-PPA 提高连接蛋白，直接增强屏障",
        "Microbiome Fig.7l-o：B. fragilis 与 3-PPA 的代谢组证据": "Microbiome Fig.7l-o：B. fragilis 与 3-PPA 形成关键菌-代谢物线索",
        "Microbiome Fig.8a-h：3-PPA / B. fragilis 激活 AhR 信号": "Microbiome Fig.8a-h：3-PPA/B. fragilis 激活 AhR 信号",
        "Microbiome Fig.8i-o：AhR 抑制剂阻断 3-PPA 的屏障保护作用": "Microbiome Fig.8i-o：AhR 抑制剂削弱 3-PPA 屏障保护",
        "Microbiome Fig.8p-s：AhR 抑制后连接蛋白改善被阻断": "Microbiome Fig.8p-s：AhR 被抑制后连接蛋白改善被阻断",
        "整合机制假说：植物酚类—菌群—代谢物—AhR/Nrf2—屏障": "整合机制假说：植物酚类—菌群—代谢物—AhR/Nrf2—屏障",
        "对荔枝酚类课题的启发": "课题启发：荔枝酚类可沿“菌群代谢物—屏障修复”验证",
        "研究总结与亮点": "研究总结：主文献证据扎实，Microbiome 提升机制深度",
    }

    for slide in prs.slides:
        old = first_title(slide)
        if old in title_updates:
            set_title(slide, title_updates[old])

    # Split Figure 2B-E.
    fig2_slide = next(s for s in prs.slides if first_title(s).startswith("Figure 2B-D"))
    set_tag(fig2_slide, "Fig.2B-D")
    remove_pictures(fig2_slide)
    fig2_slide.shapes.add_picture(str(bd_img), in_(0.82), in_(0.95), width=in_(7.25), height=in_(5.45))
    update_body(
        fig2_slide,
        "图怎么看",
        "B：villus height，代表吸收面积\nC：crypt depth，反映隐窝状态\nD：V/C ratio，综合评价肠道形态\nEA 提高绒毛高度和 V/C，提示吸收结构改善",
    )
    fig2e_slide = add_figure_slide(
        prs,
        "Figure 2E：EA 降低 TUNEL AOD，提示上皮凋亡减少",
        "Fig.2E",
        [e_img],
        "本页结论",
        "TUNEL staining AOD 越高，说明凋亡信号越强。\nEA 降低 TUNEL AOD，FEA 也呈下降趋势。\n这说明 EA 可减轻断奶应激导致的上皮细胞损伤。",
    )
    fig2_index = list(prs.slides).index(fig2_slide)
    move_last_slide_to(prs, fig2_index + 1)

    # Split Microbiome Fig.2e-k.
    micro_slide = next(s for s in prs.slides if first_title(s).startswith("Microbiome Fig.2e-g"))
    set_tag(micro_slide, "Micro Fig.2")
    pics = [s for s in micro_slide.shapes if s.shape_type == 13]
    blobs = []
    for p in pics:
        blobs.append((p.image.blob, p.image.ext))
    # Keep only e-g picture on this slide and make it larger.
    for p in pics[1:]:
        remove_shape(p)
    if pics:
        pics[0].left = in_(0.75)
        pics[0].top = in_(1.25)
        pics[0].width = in_(7.45)
        pics[0].height = in_(2.55)
    update_body(
        micro_slide,
        "结果要点",
        "e-g：比较 DLY-R 与 CM-R 小鼠屏障损伤指标。\nCM-R 组 DAO、内毒素和 D-乳酸等指标下降。\n说明 CM 菌群移植可转移更强屏障表型。",
    )
    wb = ASSETS / "micro_fig2h_wb.png"
    quant = ASSETS / "micro_fig2i_k_quant.png"
    if len(blobs) >= 3:
        wb.write_bytes(blobs[1][0])
        quant.write_bytes(blobs[2][0])
    micro2_slide = add_figure_slide(
        prs,
        "Microbiome Fig.2h-k：CM 菌群提高 ZO-1、E-cadherin 和 Connexin 43",
        "Micro Fig.2",
        [wb, quant],
        "本页结论",
        "h：Western blot 检测连接蛋白。\ni-k：定量显示 CM-R 组 ZO-1、E-cadherin、Connexin 43 升高。\n说明菌群移植改善屏障结构蛋白。",
    )
    micro_index = list(prs.slides).index(micro_slide)
    move_last_slide_to(prs, micro_index + 1)

    # Repurpose old thanks slide into Day 1 reading route.
    thanks = next(s for s in prs.slides if first_title(s).startswith("致"))
    remove_content_except_footer(thanks)
    add_title_bar(thanks, "两天内读懂本次汇报文献的阅读路线：Day 1", "Reading")
    add_text_block(
        thanks,
        0.7,
        1.05,
        5.9,
        5.45,
        "Day 1：吃透 Animal Nutrition 主文献",
        [
            "1. 先读 Abstract，抓住 Ctrl、EA、FEA 三组",
            "2. 再读 Methods：126 头断奶仔猪、0.1% EA、14 d、FEA 粪菌移植",
            "3. 按 Figure 顺序读：Fig.1 表型/屏障/氧化应激；Fig.2 形态/凋亡；Fig.3 基因",
            "4. Fig.4–5 看菌群结构和差异菌；Fig.6–7 看 SCFAs 和相关性",
            "5. 最后读 Discussion，只找作者如何解释 EA—菌群—抗氧化—屏障",
        ],
        14,
    )
    add_text_block(
        thanks,
        7.0,
        1.05,
        5.3,
        5.45,
        "Day 1 目标",
        [
            "用 3 分钟讲清：",
            "• 实验设计：Ctrl / EA / FEA",
            "• 主要结果：腹泻下降、屏障改善、MDA 降低、菌群和 SCFAs 改变",
            "• 核心结论：EA 通过菌群参与改善断奶应激",
            "背熟一句话：FEA 组是全文设计亮点，用来判断 EA 作用是否由菌群介导。",
        ],
        15,
    )

    add_content_slide(
        prs,
        "Day 2：读懂 Microbiome 机制补充，并串回主文献",
        "Reading",
        [
            (
                0.7,
                1.05,
                5.9,
                5.45,
                "阅读顺序",
                [
                    "1. 先读 Abstract，抓住 pig gut microbiota、B. fragilis、3-PPA、AhR、intestinal barrier",
                    "2. 不需要全文展开，重点看 3-PPA/AhR 相关结果图",
                    "3. 只回答一个问题：菌群代谢物如何通过宿主信号修复猪肠屏障？",
                    "4. 最后串回主文献：EA 证明调菌群，Microbiome 解释代谢物—AhR—屏障",
                ],
                14,
            ),
            (
                7.0,
                1.05,
                5.3,
                5.45,
                "Day 2 目标",
                [
                    "用 2 分钟解释：",
                    "• 为什么要补充 Microbiome",
                    "• 它如何提升主文献机制深度",
                    "• 它不能替代 EA 研究，也不能证明荔枝酚类已经走 3-PPA/AhR",
                    "收束句：这是机制启发，不是主文献直接证据。",
                ],
                15,
            ),
        ],
    )

    add_content_slide(
        prs,
        "面向博士/老师的汇报策略",
        "Strategy",
        [
            (
                0.7,
                1.05,
                5.9,
                5.45,
                "讲法边界",
                [
                    "1. 不要把两篇文章说成一个研究：Animal Nutrition 是主文献，Microbiome 是机制补充",
                    "2. 不要把假说说成已验证结论：荔枝酚类通过 3-PPA/AhR 属于待验证",
                    "3. 每个 Figure 都回答一个问题：有效吗、修复吗、调菌吗、是否有代谢物线索",
                    "4. 主动讲局限：主文献没有直接做 3-PPA/AhR，Microbiome 也不是 EA 干预研究",
                ],
                14,
            ),
            (
                7.0,
                1.05,
                5.3,
                5.45,
                "落到本人课题",
                [
                    "后续可检测：",
                    "• MDA、SOD、CAT、T-AOC、GSH/GSSG",
                    "• ZO-1、Occludin、Claudin、MUC2",
                    "• 16S、SCFAs、芳香族代谢物",
                    "• AhR、Nrf2、NF-κB",
                    "最后把 EA 作为酚类模型，把荔枝酚类作为待验证对象。",
                ],
                15,
            ),
        ],
    )

    add_content_slide(
        prs,
        "可能被老师问到的问题与回答",
        "Q&A",
        [
            (
                0.55,
                0.98,
                6.1,
                5.65,
                "Q1–Q3：主线和边界",
                [
                    "Q1：为什么主讲 Animal Nutrition？",
                    "A：它最贴合断奶仔猪 + 酚类 + 氧化应激 + 屏障 + 菌群；Microbiome 只补机制。",
                    "Q2：FEA 组意义？",
                    "A：接受 EA 组粪菌；若能复制部分效果，说明 EA 保护至少部分由菌群介导。",
                    "Q3：主文献证明 3-PPA/AhR 吗？",
                    "A：没有，3-PPA/AhR 只能作为 Microbiome 的机制补充。",
                ],
                12,
            ),
            (
                6.9,
                0.98,
                5.85,
                5.65,
                "Q4–Q5：课题借鉴",
                [
                    "Q4：荔枝酚类课题借鉴什么？",
                    "A：借鉴动物模型、指标体系和 FEA 设计，并增加芳香族代谢物、AhR、Nrf2、NF-κB。",
                    "Q5：两篇文章是不是同一组？",
                    "A：不是。Animal Nutrition 是主文献，Microbiome 是独立机制补充。",
                    "一句话收束：两者共同服务于“植物酚类—菌群—代谢物—屏障”假说。",
                ],
                12,
            ),
        ],
    )

    final_thanks = add_content_slide(
        prs,
        "致 谢",
        "Thanks",
        [
            (
                2.25,
                2.25,
                8.7,
                2.55,
                "感谢各位老师批评指正",
                [
                    "主线：EA 是酚类模型；FEA 支持菌群介导；Microbiome 补充 3-PPA/AhR；荔枝酚类是后续待验证假说。",
                ],
                18,
            )
        ],
    )

    renumber(prs.slides)
    titles = [first_title(s) for s in prs.slides]
    INTERMEDIATE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(INTERMEDIATE)
    patch_notes(INTERMEDIATE, OUT, titles)
    print(OUT)
    print(len(titles))
    for i, t in enumerate(titles, 1):
        print(f"{i:02d} {t}")


if __name__ == "__main__":
    main()
